#!/usr/bin/env python
"""Build an animated PowerPoint slide of 50 random fake rainfall document images.

Identical in style to ``make_thornton_stack_pptx.py``: 50 register-style scans
each fly in from the right at full slide height. The first lands flush against
the left edge, and every later image lands the same size, offset horizontally to
the right (same vertical position), so the pile fans out across the slide.

python-pptx has no animation API, so the PowerPoint <p:timing> tree (a standard
"Fly In / From Right" entrance per picture, each starting *After Previous*) is
built as XML and grafted onto the slide. Playback is hands-free.

Run in the ADRQ environment. Requires python-pptx:

    conda activate ADRQ
    pip install python-pptx
    python scripts/make_fake_stack_pptx.py

Output: $PDIR/fake_document_stack.pptx
"""

from __future__ import annotations

import os
import random
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml


# --------------------------------------------------------------------------- #
# Configuration
# --------------------------------------------------------------------------- #
IMAGE_ROOT = Path(
    os.environ.get(
        "FAKE_IMAGE_ROOT",
        str(Path.home() / "Projects/Auto-Daily-Rainfall-MO/fake_daily_rainfall/images"),
    )
)
OUTPUT_PATH = Path(os.environ.get("PDIR", "/data/scratch/philip.brohan/ADRQ")) / \
    "fake_document_stack.pptx"

N_IMAGES = 50           # how many random scans to draw
RANDOM_SEED = 20260811  # fixed for reproducibility (set to None for fresh picks)

# Slide geometry (16:9 widescreen).
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Every scan is drawn full slide height; width follows the image's own aspect
# ratio. The first image sits against the left edge and each later image is
# offset horizontally to the right (same vertical position), so the pile fans
# out across the slide.
IMG_HEIGHT = SLIDE_H
MARGIN = Inches(0.15)

# Animation timing (milliseconds).
FLY_DURATION_MS = 500   # duration of the FIRST image's slide-in (slowest)
SPEED_RAMP = 5.0        # last image flies in this many times faster than the first
GAP_MS = 0              # extra pause between images (0 => back-to-back)

BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)


# --------------------------------------------------------------------------- #
# Select the image files (random sample from the fake image directory)
# --------------------------------------------------------------------------- #
def select_images(root: Path, n: int, seed: int | None) -> list[Path]:
    candidates = sorted(
        p for p in root.iterdir()
        if p.suffix.lower() in (".jpg", ".jpeg", ".png")
    )
    if not candidates:
        raise SystemExit(f"No images found under {root}.")
    rng = random.Random(seed)
    return rng.sample(candidates, min(n, len(candidates)))


# --------------------------------------------------------------------------- #
# Animation XML
# --------------------------------------------------------------------------- #
def _effect_xml(shape_id: int, cid: int, dur_ms: int, gap_ms: int) -> str:
    """One "Fly In / From Right" entrance, started After Previous.

    Uses five sequential cTn ids starting at ``cid``.
    """
    a, b, c, d, e = cid, cid + 1, cid + 2, cid + 3, cid + 4
    return f"""
    <p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
      <p:cTn id="{a}" fill="hold">
        <p:stCondLst><p:cond delay="{gap_ms}"/></p:stCondLst>
        <p:childTnLst>
          <p:par>
            <p:cTn id="{b}" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst>
              <p:childTnLst>
                <p:par>
                  <p:cTn id="{c}" presetID="4" presetClass="entr" presetSubtype="2"
                         fill="hold" grpId="0" nodeType="afterEffect">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                    <p:childTnLst>
                      <p:set>
                        <p:cBhvr>
                          <p:cTn id="{d}" dur="1" fill="hold"/>
                          <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                          <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                        </p:cBhvr>
                        <p:to><p:strVal val="visible"/></p:to>
                      </p:set>
                      <p:anim calcmode="lin" valueType="num">
                        <p:cBhvr additive="base">
                          <p:cTn id="{e}" dur="{dur_ms}" fill="hold"/>
                          <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                          <p:attrNameLst><p:attrName>ppt_x</p:attrName></p:attrNameLst>
                        </p:cBhvr>
                        <p:tavLst>
                          <p:tav tm="0"><p:val><p:strVal val="1+#ppt_w/2"/></p:val></p:tav>
                          <p:tav tm="100000"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav>
                        </p:tavLst>
                      </p:anim>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
          </p:par>
        </p:childTnLst>
      </p:cTn>
    </p:par>"""


def build_timing_xml(shape_ids: list[int], dur_ms: int, gap_ms: int) -> str:
    effects = []
    cid = 3  # ids 1 (tmRoot) and 2 (mainSeq) are reserved
    n = len(shape_ids)
    for i, sid in enumerate(shape_ids):
        # First effect: no leading gap; subsequent effects honour GAP_MS.
        this_gap = 0 if i == 0 else gap_ms
        # Speed ramps up steadily: the first image takes the full duration and
        # the last is SPEED_RAMP times faster (shorter duration).
        t = 0.0 if n <= 1 else i / (n - 1)
        speed = 1.0 + (SPEED_RAMP - 1.0) * t
        this_dur = max(1, int(round(dur_ms / speed)))
        effects.append(_effect_xml(sid, cid, this_dur, this_gap))
        cid += 5
    effects_xml = "".join(effects)
    return f"""<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>{effects_xml}
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""


# --------------------------------------------------------------------------- #
# Build the deck
# --------------------------------------------------------------------------- #
def main() -> None:
    images = select_images(IMAGE_ROOT, N_IMAGES, RANDOM_SEED)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # White background.
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR

    # First pass: place every picture at its resting position and record widths.
    pictures = []
    for path in images:
        pic = slide.shapes.add_picture(str(path), 0, 0, height=IMG_HEIGHT)
        pictures.append(pic)
    max_w = max(p.width for p in pictures)

    # All images share the same (full) height and top edge; only the horizontal
    # position advances. The first image is flush against the left edge, the
    # last image's right edge reaches the right edge of the slide.
    n = len(pictures)
    start_left = 0
    end_left = max(0, int(SLIDE_W - max_w))
    top = int((SLIDE_H - IMG_HEIGHT) // 2)  # 0 when images fill the height

    def lerp(a: int, b: int, t: float) -> int:
        return int(round(a + (b - a) * t))

    for i, pic in enumerate(pictures):
        t = 0.0 if n == 1 else i / (n - 1)
        pic.left = Emu(lerp(start_left, end_left, t))
        pic.top = Emu(top)

    # Caption (drawn last => in front).
    cap = slide.shapes.add_textbox(MARGIN, SLIDE_H - Inches(0.5),
                                   Inches(8), Inches(0.4))
    tf = cap.text_frame
    tf.text = "Fake daily-rainfall registers \u2014 50 random synthetic scans"
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # Graft the animation timing onto the slide.
    shape_ids = [p.shape_id for p in pictures]
    timing = parse_xml(build_timing_xml(shape_ids, FLY_DURATION_MS, GAP_MS))
    slide.element.append(timing)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"Wrote {OUTPUT_PATH}  ({len(pictures)} images)")


if __name__ == "__main__":
    main()
