#!/usr/bin/env python
"""Build a single PowerPoint slide illustrating the QC1 algorithm.

QC1 (the exact-monthly consistency check,
:mod:`rainfall_rescue_sqlite.parquet_qc_exact_monthly`) asks one question of every
exact-matched station-month:

    Does the month's transcribed daily rainfall add up to the *independently*
    digitised monthly total?

Concretely it sums the consensus daily values (the member median per day) over a
month, compares that sum with the Rainfall-Rescue monthly total, and if the two
agree to within a small tolerance (default 0.01 in) it stamps **every day** in
that month ``pass``; otherwise every day is ``fail``. It is deliberately
all-or-nothing per month.

This slide is a purely visual talking aid for describing that algorithm. It shows
two worked example months side by side -- one whose daily values agree with the
RR total (whole month passes, green ticks) and one that disagrees (whole month
fails, red crosses) -- so the check and its per-month verdict are obvious without
reading any prose.

The pictorial day grids are rendered with matplotlib; the number cards, compare
node, flow arrows and short captions are native (editable) PowerPoint shapes.

Run in the ADRQ environment (needs python-pptx and matplotlib):

    conda activate ADRQ
    python scripts/make_qc1_slide_pptx.py

Output: $PDIR/qc1_process.pptx
"""

from __future__ import annotations

import os
import tempfile
from pathlib import Path
from typing import List, Sequence

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from matplotlib.patches import Rectangle  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


# --------------------------------------------------------------------------- #
# Configuration
# --------------------------------------------------------------------------- #
OUTPUT_PATH = Path(os.environ.get("PDIR", "/data/scratch/philip.brohan/ADRQ")) / \
    "qc1_process.pptx"

# QC1's default agreement tolerance, in inches (parquet_qc_exact_monthly.py).
TOLERANCE_IN = 0.01

# Slide geometry (16:9 widescreen).
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette.
INK = RGBColor(0x22, 0x22, 0x22)
GREY_FILL = RGBColor(0xEC, 0xEF, 0xF1)
GREY_LINE = RGBColor(0x90, 0x9C, 0xA6)
BLUE_FILL = RGBColor(0xE3, 0xF0, 0xFB)
BLUE_LINE = RGBColor(0x15, 0x65, 0xC0)
BLUE_INK = RGBColor(0x0D, 0x47, 0xA1)
PASS_GREEN = RGBColor(0x2E, 0x7D, 0x32)
FAIL_RED = RGBColor(0xC6, 0x28, 0x28)
ARROW_FILL = RGBColor(0xB0, 0xBE, 0xC5)

# Two worked example months. Each is 30 daily rainfall values (inches). The card
# sums are computed from these arrays so the picture and the numbers always agree.
MONTH_PASS: List[float] = [
    0.00, 0.05, 0.12, 0.00, 0.30, 0.08,
    0.00, 0.00, 0.22, 0.14, 0.00, 0.06,
    0.18, 0.40, 0.10, 0.00, 0.02, 0.00,
    0.25, 0.00, 0.16, 0.34, 0.09, 0.00,
    0.11, 0.00, 0.07, 0.20, 0.13, 0.00,
]
MONTH_FAIL: List[float] = [
    0.00, 0.60, 0.00, 0.45, 0.10, 0.00,
    0.35, 0.00, 0.28, 0.00, 0.52, 0.14,
    0.00, 0.22, 0.00, 0.48, 0.06, 0.30,
    0.00, 0.19, 0.41, 0.00, 0.13, 0.00,
    0.27, 0.00, 0.38, 0.00, 0.16, 0.09,
]


# matplotlib colours (tuples, to match the pptx palette above).
MPL_INK = (0.13, 0.13, 0.13)
MPL_BLUE = (0.08, 0.40, 0.75)
MPL_GREEN = (0.18, 0.49, 0.20)
MPL_RED = (0.78, 0.16, 0.16)

MONTH_INITIALS = ["J", "F", "M", "A", "M", "J", "J", "A", "S", "O", "N", "D"]
N_MONTHS = 12
N_DAYS = 31


# --------------------------------------------------------------------------- #
# matplotlib year-grid panels
# --------------------------------------------------------------------------- #
def _draw_tick(ax, x: float, y: float, w: float) -> None:
    ax.plot(
        [x + 0.24 * w, x + 0.42 * w, x + 0.74 * w],
        [y + 0.58, y + 0.74, y + 0.28],
        color=MPL_GREEN, lw=1.7, solid_capstyle="round", solid_joinstyle="round",
    )


def _draw_cross(ax, x: float, y: float, w: float) -> None:
    ax.plot([x + 0.30 * w, x + 0.70 * w], [y + 0.30, y + 0.70],
            color=MPL_RED, lw=1.7, solid_capstyle="round")
    ax.plot([x + 0.30 * w, x + 0.70 * w], [y + 0.70, y + 0.30],
            color=MPL_RED, lw=1.7, solid_capstyle="round")


def render_year_grid(values: Sequence[float], mode: str, out_path: Path,
                     *, highlight_col: int) -> None:
    """Render a station-year as a 12-month x 31-day grid, one month highlighted.

    The grid mimics the real data layout: 12 columns (months) by 31 rows (days).
    The month being checked is drawn as a widened, outlined column; only that
    column carries content (the day values in ``"plain"`` mode, or green ticks /
    red crosses in ``"pass"`` / ``"fail"`` mode). In ``"plain"`` mode a bracket and
    a Sigma below the column show that the monthly total is that column's sum.
    """
    accent = {"plain": MPL_BLUE, "pass": MPL_GREEN, "fail": MPL_RED}[mode]

    # Variable column widths: the highlighted month is wider so its numbers read.
    w_norm, w_hi = 1.0, 2.6
    col_w = [w_hi if i == highlight_col else w_norm for i in range(N_MONTHS)]
    xs = [0.0]
    for w in col_w:
        xs.append(xs[-1] + w)
    total_w = xs[-1]

    # Fixed margins (same for every mode) so grids line up across panels: room
    # above for month labels, below for the column-sum bracket.
    top_pad, bot_pad = 1.6, 2.2
    fig, ax = plt.subplots(figsize=(2.35, 2.9), dpi=220)
    ax.set_xlim(0, total_w)
    ax.set_ylim(-top_pad, N_DAYS + bot_pad)
    ax.invert_yaxis()
    ax.axis("off")

    # Month initials along the top; the checked month stands out.
    for i, label in enumerate(MONTH_INITIALS):
        is_hi = i == highlight_col
        ax.text(
            xs[i] + col_w[i] / 2, -0.75, label, ha="center", va="center",
            fontsize=8.0 if is_hi else 6.0,
            color=accent if is_hi else (0.55, 0.55, 0.55),
            fontweight="bold" if is_hi else "normal",
        )

    # Cells.
    for i in range(N_MONTHS):
        is_hi = i == highlight_col
        for d in range(N_DAYS):
            x = xs[i]
            has_val = is_hi and d < len(values)
            if is_hi:
                if d >= len(values):
                    face = (0.96, 0.96, 0.96)
                elif mode == "plain":
                    face = "white"
                elif mode == "pass":
                    face = (0.90, 0.96, 0.90)
                else:
                    face = (0.99, 0.92, 0.92)
            else:
                face = (0.945, 0.95, 0.955)
            ax.add_patch(
                Rectangle((x, d), col_w[i], 1, facecolor=face,
                          edgecolor="0.80", linewidth=0.4)
            )
            if has_val:
                if mode == "plain":
                    ax.text(x + col_w[i] / 2, d + 0.5, f"{values[d]:.2f}",
                            ha="center", va="center", fontsize=5.3, color="0.12")
                elif mode == "pass":
                    _draw_tick(ax, x, d, col_w[i])
                else:
                    _draw_cross(ax, x, d, col_w[i])

    # Emphasise the checked column.
    xh = xs[highlight_col]
    ax.add_patch(
        Rectangle((xh, 0), w_hi, N_DAYS, fill=False, edgecolor=accent, linewidth=1.8)
    )

    # Column-sum cue below the highlighted month (numbers view only).
    if mode == "plain":
        xr = xh + w_hi
        yb = N_DAYS + 0.5
        ax.plot([xh, xh, xr, xr], [yb - 0.28, yb, yb, yb - 0.28],
                color=MPL_INK, lw=1.2, solid_capstyle="round")
        ax.annotate("", xy=((xh + xr) / 2, yb + 1.55), xytext=((xh + xr) / 2, yb + 0.15),
                    arrowprops=dict(arrowstyle="-|>", color=MPL_INK, lw=1.5))
        ax.text((xh + xr) / 2, yb + 1.35, "\u03a3", ha="center", va="center",
                fontsize=11, color=MPL_INK, fontweight="bold")

    fig.subplots_adjust(left=0.01, right=0.99, top=0.99, bottom=0.01)
    fig.savefig(out_path, transparent=True, pad_inches=0.02)
    plt.close(fig)


# --------------------------------------------------------------------------- #
# python-pptx helpers
# --------------------------------------------------------------------------- #
def _style_text(shape, text: str, *, size: int, color: RGBColor, bold: bool) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_card(slide, left, top, width, height, text, *, fill, line, ink,
             size=18, bold=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    _style_text(shp, text, size=size, color=ink, bold=bold)
    return shp


def add_arrow(slide, left, top, width, height):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = ARROW_FILL
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_caption(slide, left, top, width, text, *, size=12, color=INK, bold=False):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(0.32)
    )
    _style_text(box, text, size=size, color=color, bold=bold)
    return box


def add_picture(slide, path, left, top, width, height):
    return slide.shapes.add_picture(
        str(path), Inches(left), Inches(top), Inches(width), Inches(height)
    )


# --------------------------------------------------------------------------- #
# Row layout
# --------------------------------------------------------------------------- #
def build_row(slide, tmp: Path, *, row_top: float, values: Sequence[float],
              verdict: str, highlight_col: int) -> None:
    """Lay out one worked-example month across the slide.

    ``verdict`` is ``"pass"`` or ``"fail"``. ``highlight_col`` is the month column
    (0-based) that is being checked. The RR total is set equal to the daily sum
    for a pass, or clearly different for a fail, so the compare node tells the
    true story.
    """
    total = sum(values)
    passed = verdict == "pass"
    rr_total = total if passed else 2.10  # a wrong monthly figure for the fail case
    result_ink = PASS_GREEN if passed else FAIL_RED

    grid_w, grid_h = 2.35, 2.86
    mid = row_top + grid_h / 2.0

    # 1. Year of daily values as a 12-month x 31-day grid (picture).
    plain = tmp / f"grid_plain_{verdict}.png"
    render_year_grid(values, "plain", plain, highlight_col=highlight_col)
    add_picture(slide, plain, 0.35, row_top, grid_w, grid_h)
    add_caption(slide, 0.20, row_top + grid_h + 0.00, grid_w + 0.3,
                "Daily values \u2014 one month per column", size=12, bold=True)

    # 2. Column-sum arrow.
    add_arrow(slide, 2.95, mid - 0.22, 0.70, 0.44)
    add_caption(slide, 2.75, mid - 0.66, 1.10, "Column sum", size=11.5)

    # 3. Number cards: computed monthly total and the independent RR total.
    card_w, card_h = 2.15, 0.80
    add_caption(slide, 3.85, mid - 1.36, card_w, "Monthly total (\u03a3 daily)",
                size=11.5, bold=True)
    add_card(slide, 3.85, mid - 1.04, card_w, card_h,
             f"\u03a3 = {total:.2f} in", fill=GREY_FILL, line=GREY_LINE, ink=INK,
             size=17)
    add_card(slide, 3.85, mid + 0.24, card_w, card_h,
             f"{rr_total:.2f} in", fill=BLUE_FILL, line=BLUE_LINE, ink=BLUE_INK,
             size=17)
    add_caption(slide, 3.85, mid + 1.06, card_w, "RR total (independent)",
                size=11.5, color=BLUE_INK, bold=True)

    # 4. Compare node.
    add_card(slide, 6.40, mid - 0.62, 1.55, 1.24,
             f"|\u03a3 \u2212 RR|\n\u2264 {TOLERANCE_IN:.2f} in ?",
             fill=RGBColor(0xFF, 0xFF, 0xFF), line=INK, ink=INK, size=13,
             shape_type=MSO_SHAPE.DIAMOND)

    # 5. Verdict arrow.
    add_arrow(slide, 8.15, mid - 0.22, 0.70, 0.44)

    # 6. Verdict grid (picture) + caption.
    verdict_img = tmp / f"grid_{verdict}.png"
    render_year_grid(values, verdict, verdict_img, highlight_col=highlight_col)
    add_picture(slide, verdict_img, 9.05, row_top, grid_w, grid_h)
    label = "Pass \u2014 whole month" if passed else "Fail \u2014 whole month"
    add_caption(slide, 8.90, row_top + grid_h + 0.00, grid_w + 0.3, label,
                size=13, color=result_ink, bold=True)


# --------------------------------------------------------------------------- #
# Slide assembly
# --------------------------------------------------------------------------- #
def build_slide(output_path: Path = OUTPUT_PATH) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    title = slide.shapes.add_textbox(Inches(0), Inches(0.12), SLIDE_W, Inches(0.5))
    _style_text(title, "QC1 \u2014 monthly-total consistency check",
                size=23, color=INK, bold=True)

    with tempfile.TemporaryDirectory(dir="/var/tmp") as td:
        tmp = Path(td)
        # Two worked example months in the same year grid: a different column is
        # checked each time to underline that every month is judged on its own.
        build_row(slide, tmp, row_top=0.68, values=MONTH_PASS, verdict="pass",
                  highlight_col=3)
        build_row(slide, tmp, row_top=3.92, values=MONTH_FAIL, verdict="fail",
                  highlight_col=8)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

    return output_path


def main() -> None:
    path = build_slide()
    print(f"Wrote QC1 process slide -> {path}")


if __name__ == "__main__":
    main()
