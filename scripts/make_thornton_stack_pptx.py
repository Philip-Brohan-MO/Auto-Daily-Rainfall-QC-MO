#!/usr/bin/env python
"""Build an animated PowerPoint slide of the Thornton Reservoir rainfall images.

The 90 annual register scans (1871-1960) each fly in from the right at full
slide height. The first year lands flush against the left edge, and every later
year lands the same size, offset horizontally to the right (same vertical
position), so the pile fans out across the slide as the years advance.

python-pptx has no animation API, so the PowerPoint <p:timing> tree (a standard
"Fly In / From Right" entrance per picture, each starting *After Previous*) is
built as XML and grafted onto the slide. Playback is hands-free.

Run in the ADRQ environment. Requires python-pptx:

    conda activate ADRQ
    pip install python-pptx
    python scripts/make_thornton_stack_pptx.py

Output: $PDIR/thornton_reservoir_stack.pptx
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml


# --------------------------------------------------------------------------- #
# Configuration
# --------------------------------------------------------------------------- #
IMAGE_ROOT = Path(
    os.environ.get(
        "DRAIN_IMAGE_ROOT",
        "/data/scratch/philip.brohan/documents/Daily_Rainfall_UK/jpgs_25pc_filtered",
    )
)
OUTPUT_PATH = Path(os.environ.get("PDIR", "/data/scratch/philip.brohan/ADRQ")) / \
    "thornton_reservoir_stack.pptx"

# Slide geometry (16:9 widescreen).
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Every scan is drawn full slide height; width follows the image's own aspect
# ratio. The first image sits against the left edge and each later year is
# offset horizontally to the right (same vertical position), so the pile fans
# out across the slide as the years advance.
IMG_HEIGHT = SLIDE_H
MARGIN = Inches(0.15)

# Animation timing (milliseconds).
FLY_DURATION_MS = 500   # duration of the FIRST image's slide-in (slowest)
SPEED_RAMP = 5.0        # last image flies in this many times faster than the first
GAP_MS = 0              # extra pause between images (0 => back-to-back)

BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)

# The 90 image identifiers, one per year 1871-1960 (in chronological order).
IDENTIFIERS = [
    "DRain_1871-1880_Leicestershire-105",  # 1871
    "DRain_1871-1880_Leicestershire-107",  # 1872
    "DRain_1871-1880_Leicestershire-108",  # 1873
    "DRain_1871-1880_Leicestershire-109",  # 1874
    "DRain_1871-1880_Leicestershire-110",  # 1875
    "DRain_1871-1880_Leicestershire-111",  # 1876
    "DRain_1871-1880_Leicestershire-112",  # 1877
    "DRain_1871-1880_Leicestershire-113",  # 1878
    "DRain_1871-1880_Leicestershire-114",  # 1879
    "DRain_1871-1880_Leicestershire-115",  # 1880
    "DRain_1881-1890_Leicestershire-151",  # 1881
    "DRain_1881-1890_Leicestershire-152",  # 1882
    "DRain_1881-1890_Leicestershire-153",  # 1883
    "DRain_1881-1890_Leicestershire-154",  # 1884
    "DRain_1881-1890_Leicestershire-155",  # 1885
    "DRain_1881-1890_Leicestershire-156",  # 1886
    "DRain_1881-1890_Leicestershire-157",  # 1887
    "DRain_1881-1890_Leicestershire-158",  # 1888
    "DRain_1881-1890_Leicestershire-159",  # 1889
    "DRain_1881-1890_Leicestershire-160",  # 1890
    "DRain_1891-1900_RainNos_Leicestershire_Lincolnshire_B024-1",   # 1891
    "DRain_1891-1900_RainNos_Leicestershire_Lincolnshire_B024-3",   # 1892
    "DRain_1891-1900_RainNos_Leicestershire_Lincolnshire_B024-4",   # 1893
    "DRain_1891-1900_RainNos_Leicestershire_Lincolnshire_B024-5",   # 1894
    "DRain_1891-1900_RainNos_Leicestershire_Lincolnshire_B024-6",   # 1895
    "DRain_1891-1900_RainNos_Leicestershire_Lincolnshire_B024-9",   # 1896
    "DRain_1891-1900_RainNos_Leicestershire_Lincolnshire_B024-10",  # 1897
    "DRain_1891-1900_RainNos_Leicestershire_Lincolnshire_B024-11",  # 1898
    "DRain_1891-1900_RainNos_Leicestershire_Lincolnshire_B024-12",  # 1899
    "DRain_1891-1900_RainNos_Leicestershire_Lincolnshire_B024-13",  # 1900
    "DRain_1901-1910_RainNos_Leicestershire_B031-1",   # 1901
    "DRain_1901-1910_RainNos_Leicestershire-237",      # 1902
    "DRain_1901-1910_RainNos_Leicestershire_B031-3",   # 1903
    "DRain_1901-1910_RainNos_Leicestershire_B031-4",   # 1904
    "DRain_1901-1910_RainNos_Leicestershire_B031-5",   # 1905
    "DRain_1901-1910_RainNos_Leicestershire_B031-6",   # 1906
    "DRain_1901-1910_RainNos_Leicestershire_B031-7",   # 1907
    "DRain_1901-1910_RainNos_Leicestershire_B031-8",   # 1908
    "DRain_1901-1910_RainNos_Leicestershire-244",      # 1909
    "DRain_1901-1910_RainNos_Leicestershire-245",      # 1910
    "DRain_1911-1920_RainNos_Lancashire_V-W_Leceistershire_B063-0",   # 1911
    "DRain_1911-1920_RainNos_Lancashire_V-W_Leceistershire_B063-2",   # 1912
    "DRain_1911-1920_RainNos_Lancashire_V-W_Leceistershire_B063-3",   # 1913
    "DRain_1911-1920_RainNos_Lancashire_V-W_Leceistershire_B063-4",   # 1914
    "DRain_1911-1920_RainNos_Lancashire_V-W_Leceistershire_B063-5",   # 1915
    "DRain_1911-1920_RainNos_Lancashire_V-W_Leceistershire_B063-6",   # 1916
    "DRain_1911-1920_RainNos_Lancashire_V-W_Leceistershire_B063-7",   # 1917
    "DRain_1911-1920_RainNos_Lancashire_V-W_Leceistershire_B063-8",   # 1918
    "DRain_1911-1920_RainNos_Lancashire_V-W_Leceistershire_B063-9",   # 1919
    "DRain_1911-1920_RainNos_Lancashire_V-W_Leceistershire_B063-10",  # 1920
    "DRain_1921-1930_RainNos_Leicestershire_B046-0",   # 1921
    "DRain_1921-1930_RainNos_Leicestershire_B046-2",   # 1922
    "DRain_1921-1930_RainNos_Leicestershire_B046-3",   # 1923
    "DRain_1921-1930_RainNos_Leicestershire_B046-4",   # 1924
    "DRain_1921-1930_RainNos_Leicestershire_B046-5",   # 1925
    "DRain_1921-1930_RainNos_Leicestershire_B046-6",   # 1926
    "DRain_1921-1930_RainNos_Leicestershire_B046-7",   # 1927
    "DRain_1921-1930_RainNos_Leicestershire_B046-8",   # 1928
    "DRain_1921-1930_RainNos_Leicestershire_B046-9",   # 1929
    "DRain_1921-1930_RainNos_Leicestershire_B046-10",  # 1930
    "DRain_1931-1940_RainNos_2489-2516-293",  # 1931
    "DRain_1931-1940_RainNos_2489-2516-295",  # 1932
    "DRain_1931-1940_RainNos_2489-2516-296",  # 1933
    "DRain_1931-1940_RainNos_2489-2516-297",  # 1934
    "DRain_1931-1940_RainNos_2489-2516-298",  # 1935
    "DRain_1931-1940_RainNos_2489-2516-299",  # 1936
    "DRain_1931-1940_RainNos_2489-2516-300",  # 1937
    "DRain_1931-1940_RainNos_2489-2516-302",  # 1938
    "DRain_1931-1940_RainNos_2489-2516-303",  # 1939
    "DRain_1931-1940_RainNos_2489-2516_B035-11",  # 1940
    "DRain_1941-1950_RainNos_2494-2525_B021-1",   # 1941
    "DRain_1941-1950_RainNos_2494-2525-196",      # 1942
    "DRain_1941-1950_RainNos_2494-2525_B021-3",   # 1943
    "DRain_1941-1950_RainNos_2494-2525_B021-4",   # 1944
    "DRain_1941-1950_RainNos_2494-2525_B021-5",   # 1945
    "DRain_1941-1950_RainNos_2494-2525_B021-6",   # 1946
    "DRain_1941-1950_RainNos_2494-2525_B021-7",   # 1947
    "DRain_1941-1950_RainNos_2494-2525_B021-8",   # 1948
    "DRain_1941-1950_RainNos_2494-2525-204",      # 1949
    "DRain_1941-1950_RainNos_2494-2525_B021-12",  # 1950
    "DRain_1951-1962_RainNos_2501-2526_B011-0",   # 1951
    "DRain_1951-1962_RainNos_2501-2526_B011-3",   # 1952
    "DRain_1951-1962_RainNos_2501-2526_B011-4",   # 1953
    "DRain_1951-1962_RainNos_2501-2526_B011-5",   # 1954
    "DRain_1951-1962_RainNos_2501-2526_B011-8",   # 1955
    "DRain_1951-1962_RainNos_2501-2526_B011-9",   # 1956
    "DRain_1951-1962_RainNos_2501-2526_B011-12",  # 1957
    "DRain_1951-1962_RainNos_2501-2526-117",      # 1958
    "DRain_1951-1962_RainNos_2501-2526_B011-15",  # 1959
    "DRain_1951-1962_RainNos_2501-2526_B011-17",  # 1960
]


# --------------------------------------------------------------------------- #
# Locate the image files (one walk of the tree -> stem -> path map)
# --------------------------------------------------------------------------- #
def build_image_index(root: Path) -> dict[str, Path]:
    index: dict[str, Path] = {}
    for dirpath, _dirs, files in os.walk(root):
        for fn in files:
            if fn.lower().endswith((".jpg", ".jpeg", ".png")):
                index.setdefault(Path(fn).stem, Path(dirpath) / fn)
    return index


# --------------------------------------------------------------------------- #
# Animation XML
# --------------------------------------------------------------------------- #
def _effect_xml(shape_id: int, cid: int, dur_ms: int, gap_ms: int) -> str:
    """One "Fly In / From Right" entrance, started After Previous.

    Uses five sequential cTn ids starting at ``cid``.
    """
    a, b, c, d, e = cid, cid + 1, cid + 2, cid + 3, cid + 4
    return f"""
    <p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
      <p:cTn id="{a}" fill="hold">
        <p:stCondLst><p:cond delay="{gap_ms}"/></p:stCondLst>
        <p:childTnLst>
          <p:par>
            <p:cTn id="{b}" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst>
              <p:childTnLst>
                <p:par>
                  <p:cTn id="{c}" presetID="4" presetClass="entr" presetSubtype="2"
                         fill="hold" grpId="0" nodeType="afterEffect">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                    <p:childTnLst>
                      <p:set>
                        <p:cBhvr>
                          <p:cTn id="{d}" dur="1" fill="hold"/>
                          <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                          <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                        </p:cBhvr>
                        <p:to><p:strVal val="visible"/></p:to>
                      </p:set>
                      <p:anim calcmode="lin" valueType="num">
                        <p:cBhvr additive="base">
                          <p:cTn id="{e}" dur="{dur_ms}" fill="hold"/>
                          <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                          <p:attrNameLst><p:attrName>ppt_x</p:attrName></p:attrNameLst>
                        </p:cBhvr>
                        <p:tavLst>
                          <p:tav tm="0"><p:val><p:strVal val="1+#ppt_w/2"/></p:val></p:tav>
                          <p:tav tm="100000"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav>
                        </p:tavLst>
                      </p:anim>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
          </p:par>
        </p:childTnLst>
      </p:cTn>
    </p:par>"""


def build_timing_xml(shape_ids: list[int], dur_ms: int, gap_ms: int) -> str:
    effects = []
    cid = 3  # ids 1 (tmRoot) and 2 (mainSeq) are reserved
    n = len(shape_ids)
    for i, sid in enumerate(shape_ids):
        # First effect: no leading gap; subsequent effects honour GAP_MS.
        this_gap = 0 if i == 0 else gap_ms
        # Speed ramps up steadily: the first image takes the full duration and
        # the last is SPEED_RAMP times faster (shorter duration).
        t = 0.0 if n <= 1 else i / (n - 1)
        speed = 1.0 + (SPEED_RAMP - 1.0) * t
        this_dur = max(1, int(round(dur_ms / speed)))
        effects.append(_effect_xml(sid, cid, this_dur, this_gap))
        cid += 5
    effects_xml = "".join(effects)
    return f"""<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>{effects_xml}
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""


# --------------------------------------------------------------------------- #
# Build the deck
# --------------------------------------------------------------------------- #
def main() -> None:
    index = build_image_index(IMAGE_ROOT)

    resolved: list[tuple[str, Path]] = []
    missing: list[str] = []
    for spec in IDENTIFIERS:
        path = index.get(spec)
        if path is None:
            missing.append(spec)
        else:
            resolved.append((spec, path))
    if missing:
        print(f"WARNING: {len(missing)} image(s) not found under {IMAGE_ROOT}:")
        for m in missing:
            print(f"  - {m}")
    if not resolved:
        raise SystemExit("No images found; nothing to build.")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # White background.
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR

    # First pass: place every picture at its resting position and record widths.
    pictures = []
    for _spec, path in resolved:
        pic = slide.shapes.add_picture(str(path), 0, 0, height=IMG_HEIGHT)
        pictures.append(pic)
    max_w = max(p.width for p in pictures)

    # All images share the same (full) height and top edge; only the horizontal
    # position advances. The first image is flush against the left edge, the
    # last image's right edge reaches the right edge of the slide.
    n = len(pictures)
    start_left = 0
    end_left = max(0, int(SLIDE_W - max_w))
    top = int((SLIDE_H - IMG_HEIGHT) // 2)  # 0 when images fill the height

    def lerp(a: int, b: int, t: float) -> int:
        return int(round(a + (b - a) * t))

    for i, pic in enumerate(pictures):
        t = 0.0 if n == 1 else i / (n - 1)
        pic.left = Emu(lerp(start_left, end_left, t))
        pic.top = Emu(top)

    # Caption (drawn last => in front, sits in the empty lower-left corner).
    cap = slide.shapes.add_textbox(MARGIN, SLIDE_H - Inches(0.5),
                                   Inches(7), Inches(0.4))
    tf = cap.text_frame
    tf.text = "Thornton Reservoir \u2014 annual daily-rainfall registers, 1871\u20131960"
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # Graft the animation timing onto the slide.
    shape_ids = [p.shape_id for p in pictures]
    timing = parse_xml(build_timing_xml(shape_ids, FLY_DURATION_MS, GAP_MS))
    slide.element.append(timing)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"Wrote {OUTPUT_PATH}  ({len(pictures)} images"
          f"{f', {len(missing)} missing' if missing else ''})")


if __name__ == "__main__":
    main()
