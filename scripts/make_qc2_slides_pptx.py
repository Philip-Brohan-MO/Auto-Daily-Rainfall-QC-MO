#!/usr/bin/env python
"""Build a three-slide PowerPoint deck illustrating the QC2 algorithm.

QC2 (the second quality-control check) is a spatial second opinion that
re-examines the observations QC1 rejected. It runs in two stages
(:mod:`rainfall_rescue_sqlite.parquet_regional_stats` and
:mod:`rainfall_rescue_sqlite.parquet_secondary_qc`):

1. **Regional neighbour statistics.** For every located station-day, compute the
   median, count and MAD of the *QC1-pass* neighbours' consensus rainfall for the
   same calendar day, at 20 km and 50 km (the station itself excluded).
2. **Expectation models.** Two XGBoost models are trained on the reliable
   (QC1-pass) rows: Model 1 predicts a station's own rainfall from its regional
   statistics; Model 2 predicts the absolute error of Model 1. A multiplier ``k``
   is calibrated so 99% of reliable days fall inside ``pred +/- k * error``.
3. **Flagging.** Each QC1-*fail* day is re-judged against its own expectation
   range: actual inside -> ``pass`` (rescued), outside -> ``fail`` (suspect),
   no neighbours -> ``indeterminate``.

This deck is a picture-led talking aid with one slide per idea. The schematic
panels (neighbour map, calibration scatter, outcome number-lines, training table)
are rendered with matplotlib; cards, arrows and short captions are native
(editable) PowerPoint shapes. All example numbers are illustrative.

Run in the ADRQ environment (needs python-pptx and matplotlib):

    conda activate ADRQ
    python scripts/make_qc2_slides_pptx.py

Output: $PDIR/qc2_process.pptx
"""

from __future__ import annotations

import os
import tempfile
from pathlib import Path
from typing import List, Sequence, Tuple

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
import numpy as np  # noqa: E402
from matplotlib.patches import Circle, Rectangle  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


# --------------------------------------------------------------------------- #
# Configuration
# --------------------------------------------------------------------------- #
OUTPUT_PATH = Path(os.environ.get("PDIR", "/data/scratch/philip.brohan/ADRQ")) / \
    "qc2_process.pptx"

# QC2 constants (parquet_regional_stats.py / parquet_secondary_qc.py).
RADIUS_SMALL_KM = 20.0
RADIUS_LARGE_KM = 50.0
COVERAGE_TARGET = 0.99

# Slide geometry (16:9 widescreen).
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette (shared with make_qc1_slide_pptx.py, plus verdict colours).
INK = RGBColor(0x22, 0x22, 0x22)
GREY_FILL = RGBColor(0xEC, 0xEF, 0xF1)
GREY_LINE = RGBColor(0x90, 0x9C, 0xA6)
BLUE_FILL = RGBColor(0xE3, 0xF0, 0xFB)
BLUE_LINE = RGBColor(0x15, 0x65, 0xC0)
BLUE_INK = RGBColor(0x0D, 0x47, 0xA1)
GREEN_FILL = RGBColor(0xE6, 0xF4, 0xEA)
GREEN_LINE = RGBColor(0x2E, 0x7D, 0x32)
GREEN_INK = RGBColor(0x1B, 0x5E, 0x20)
RED_FILL = RGBColor(0xFB, 0xE9, 0xE7)
RED_LINE = RGBColor(0xC6, 0x28, 0x28)
RED_INK = RGBColor(0xB7, 0x1C, 0x1C)
ARROW_FILL = RGBColor(0xB0, 0xBE, 0xC5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# matplotlib colours.
MPL_INK = (0.13, 0.13, 0.13)
MPL_BLUE = (0.08, 0.40, 0.75)
MPL_GREEN = (0.18, 0.49, 0.20)
MPL_RED = (0.78, 0.16, 0.16)
MPL_GREY = (0.55, 0.55, 0.55)
MPL_AMBER = (0.85, 0.60, 0.10)


# --------------------------------------------------------------------------- #
# matplotlib panels
# --------------------------------------------------------------------------- #
def render_neighbour_map(out_path: Path) -> None:
    """A target station with its 20 km / 50 km neighbour rings.

    QC1-pass neighbours (counted) are green with their rainfall value; excluded
    neighbours (failed QC1, or beyond 50 km) are hollow grey. The target star is
    excluded from its own statistics.
    """
    fig, ax = plt.subplots(figsize=(5.3, 5.3), dpi=200)
    ax.set_aspect("equal")
    ax.axis("off")
    lim = 62
    ax.set_xlim(-lim, lim)
    ax.set_ylim(-lim, lim)

    for r in (RADIUS_SMALL_KM, RADIUS_LARGE_KM):
        ax.add_patch(Circle((0, 0), r, fill=False, ls="--",
                            ec=(0.45, 0.45, 0.45), lw=1.3))
        ax.text(0, r, f"{int(r)} km", ha="center", va="bottom", fontsize=9.5,
                color=(0.4, 0.4, 0.4),
                bbox=dict(boxstyle="round,pad=0.1", fc="white", ec="none"))

    # (x_km, y_km, passed_qc1, value_in). Distances chosen to populate both rings.
    neighbours: List[Tuple[float, float, bool, float]] = [
        (10, 6, True, 0.22), (-8, 12, True, 0.26), (14, -9, True, 0.19),
        (-13, -7, True, 0.25), (5, -15, False, 0.00), (-4, 17, True, 0.28),
        (30, 14, True, 0.20), (-28, 20, True, 0.24), (22, -30, False, 0.00),
        (-34, -18, True, 0.18), (40, -8, True, 0.23), (-44, 8, True, 0.21),
        (18, 40, True, 0.17), (-20, -40, False, 0.00), (46, 26, True, 0.15),
        (55, -34, False, 0.00), (-52, -40, False, 0.00), (8, 56, False, 0.00),
    ]
    for x, y, passed, val in neighbours:
        if passed:
            ax.scatter([x], [y], s=120, color=MPL_GREEN, edgecolor="white",
                       linewidths=0.8, zorder=4)
            ax.text(x, y - 5.5, f"{val:.2f}", ha="center", va="top", fontsize=7.5,
                    color=MPL_GREEN)
        else:
            ax.scatter([x], [y], s=90, facecolor="none", edgecolor=MPL_GREY,
                       linewidths=1.3, zorder=3)

    ax.scatter([0], [0], marker="*", s=520, color=MPL_AMBER, edgecolor="black",
               linewidths=1.0, zorder=6)
    ax.text(0, -8, "target", ha="center", va="top", fontsize=9, color=MPL_INK,
            fontweight="bold")

    # Legend.
    ax.scatter([-lim + 8], [-lim + 10], s=90, color=MPL_GREEN, edgecolor="white")
    ax.text(-lim + 13, -lim + 10, "passed QC1", va="center", fontsize=8.5,
            color=MPL_INK)
    ax.scatter([-lim + 8], [-lim + 2], s=70, facecolor="none", edgecolor=MPL_GREY,
               linewidths=1.3)
    ax.text(-lim + 13, -lim + 2, "excluded", va="center", fontsize=8.5,
            color=MPL_INK)

    fig.subplots_adjust(left=0.01, right=0.99, top=0.99, bottom=0.01)
    fig.savefig(out_path, transparent=True, pad_inches=0.02)
    plt.close(fig)


def render_training_table(out_path: Path) -> None:
    """The reliable (QC1-pass) rows: regional features -> own rainfall target."""
    headers = ["n\u2082\u2080", "med\u2082\u2080", "mad\u2082\u2080",
               "n\u2085\u2080", "med\u2085\u2080", "mad\u2085\u2080", "mon",
               "rain"]
    rows = [
        ["6", "0.24", "0.05", "18", "0.21", "0.07", "4", "0.26"],
        ["4", "0.10", "0.03", "12", "0.12", "0.04", "7", "0.09"],
        ["8", "0.55", "0.09", "21", "0.49", "0.11", "11", "0.58"],
        ["3", "0.00", "0.00", "9", "0.02", "0.02", "6", "0.00"],
        ["7", "0.31", "0.06", "16", "0.28", "0.08", "2", "0.33"],
        ["5", "0.14", "0.04", "14", "0.16", "0.05", "9", "0.12"],
    ]
    ncol = len(headers)
    nrow = len(rows) + 2  # header + data + ellipsis

    fig, ax = plt.subplots(figsize=(4.6, 4.4), dpi=200)
    ax.set_xlim(0, ncol)
    ax.set_ylim(0, nrow)
    ax.invert_yaxis()
    ax.axis("off")

    def cell(c, r, text, *, header=False, target=False):
        if header:
            face = (0.18, 0.49, 0.20)
            tcol = "white"
        elif target:
            face = (0.90, 0.96, 0.90)
            tcol = MPL_GREEN
        else:
            face = (0.965, 0.98, 0.965)
            tcol = "0.15"
        ax.add_patch(Rectangle((c, r), 1, 1, facecolor=face, edgecolor="0.8",
                               linewidth=0.6))
        ax.text(c + 0.5, r + 0.5, text, ha="center", va="center",
                fontsize=8.5 if not header else 9,
                color=tcol, fontweight="bold" if header or target else "normal")

    for c, h in enumerate(headers):
        cell(c, 0, h, header=True)
    for ri, row in enumerate(rows, start=1):
        for c, txt in enumerate(row):
            cell(c, ri, txt, target=(c == ncol - 1))
    for c in range(ncol):
        ax.text(c + 0.5, nrow - 0.5, "\u22ee", ha="center", va="center",
                fontsize=11, color="0.4")

    ax.text(ncol - 0.5, -0.9, "target", ha="center", va="center", fontsize=8,
            color=MPL_GREEN, fontweight="bold")
    ax.annotate("", xy=(ncol - 0.5, -0.25), xytext=(ncol - 0.5, -0.7),
                arrowprops=dict(arrowstyle="-|>", color=MPL_GREEN, lw=1.3))

    fig.subplots_adjust(left=0.02, right=0.98, top=0.93, bottom=0.02)
    fig.savefig(out_path, transparent=True, pad_inches=0.02)
    plt.close(fig)


def render_tree_icon(out_path: Path) -> None:
    """A schematic of a few boosted decision trees."""
    fig, ax = plt.subplots(figsize=(2.6, 1.15), dpi=200)
    ax.set_xlim(0, 3)
    ax.set_ylim(0, 1)
    ax.axis("off")
    for k in range(3):
        ox = k + 0.5
        ax.plot([ox, ox - 0.28], [0.85, 0.5], color=MPL_INK, lw=1.3)
        ax.plot([ox, ox + 0.28], [0.85, 0.5], color=MPL_INK, lw=1.3)
        ax.plot([ox - 0.28, ox - 0.42], [0.5, 0.18], color=MPL_INK, lw=1.1)
        ax.plot([ox - 0.28, ox - 0.14], [0.5, 0.18], color=MPL_INK, lw=1.1)
        ax.plot([ox + 0.28, ox + 0.14], [0.5, 0.18], color=MPL_INK, lw=1.1)
        ax.plot([ox + 0.28, ox + 0.42], [0.5, 0.18], color=MPL_INK, lw=1.1)
        ax.scatter([ox], [0.85], s=40, color=MPL_BLUE, zorder=3)
        for lx in (ox - 0.42, ox - 0.14, ox + 0.14, ox + 0.42):
            ax.scatter([lx], [0.18], s=26, color=MPL_GREEN, zorder=3)
        if k < 2:
            ax.annotate("", xy=(k + 1.02, 0.5), xytext=(k + 0.94, 0.5),
                        arrowprops=dict(arrowstyle="-|>", color=MPL_GREY, lw=1.2))
    fig.subplots_adjust(left=0.01, right=0.99, top=0.99, bottom=0.01)
    fig.savefig(out_path, transparent=True, pad_inches=0.02)
    plt.close(fig)


def render_calibration_scatter(out_path: Path) -> None:
    """Predicted vs actual on the QC1-pass calibration set, with the +/-k band."""
    rng = np.random.default_rng(11)
    n = 260
    actual = np.abs(rng.gamma(1.6, 0.14, n))
    half = 0.11  # illustrative k * error half-width
    noise = rng.normal(0, 0.05, n)
    predicted = actual + noise
    # Force ~1% of points clearly outside the band for the illustration.
    outlier_idx = rng.choice(n, size=max(2, n // 100), replace=False)
    predicted[outlier_idx] = actual[outlier_idx] + rng.choice([-1, 1],
                                                              len(outlier_idx)) * 0.22
    inside = np.abs(predicted - actual) <= half

    fig, ax = plt.subplots(figsize=(3.5, 3.5), dpi=200)
    hi = 0.9
    ax.set_xlim(0, hi)
    ax.set_ylim(0, hi)
    ax.set_aspect("equal")
    xs = np.array([0, hi])
    ax.fill_between(xs, xs - half, xs + half, color=(0.85, 0.9, 0.97), zorder=0,
                    label="pred \u00b1 k\u00b7error")
    ax.plot(xs, xs, color=MPL_BLUE, lw=1.3, zorder=1)
    ax.scatter(actual[inside], predicted[inside], s=10, color=MPL_GREY, alpha=0.7,
               zorder=2)
    ax.scatter(actual[~inside], predicted[~inside], s=16, color=MPL_RED, zorder=3)
    ax.set_xlabel("actual rainfall (in)", fontsize=8.5)
    ax.set_ylabel("predicted (in)", fontsize=8.5)
    ax.tick_params(labelsize=7.5)
    ax.text(0.05, hi - 0.06, f"\u2248{int(COVERAGE_TARGET*100)}% inside band",
            fontsize=9, color=MPL_INK, fontweight="bold", va="top")
    fig.subplots_adjust(left=0.16, right=0.97, top=0.97, bottom=0.14)
    fig.savefig(out_path, transparent=True, pad_inches=0.02)
    plt.close(fig)


def render_outcomes(out_path: Path) -> None:
    """Three QC1-fail days re-judged against their expectation range."""
    fig, axes = plt.subplots(3, 1, figsize=(9.6, 2.7), dpi=200, sharex=True)
    xlo, xhi = 0.0, 0.8
    pred, half = 0.28, 0.12
    cases = [
        ("inside \u2192 pass (rescued)", MPL_GREEN, 0.31, True),
        ("outside \u2192 fail (suspect)", MPL_RED, 0.62, True),
        ("no neighbours \u2192 indeterminate", MPL_GREY, 0.40, False),
    ]
    for ax, (label, color, actual, has_band) in zip(axes, cases):
        ax.set_xlim(xlo, xhi)
        ax.set_ylim(0, 1)
        ax.set_yticks([])
        for spine in ("left", "right", "top"):
            ax.spines[spine].set_visible(False)
        ax.spines["bottom"].set_color("0.6")
        ax.set_ylabel(label, rotation=0, ha="right", va="center", fontsize=9,
                      color=color, fontweight="bold", labelpad=8)
        if has_band:
            ax.axvspan(pred - half, pred + half, ymin=0.30, ymax=0.70,
                       color=(0.85, 0.9, 0.97), zorder=0)
            ax.plot([pred, pred], [0.28, 0.72], color=MPL_BLUE, lw=1.4, zorder=1)
            ax.scatter([actual], [0.5], marker="D", s=95, color=color,
                       edgecolor="black", linewidths=0.5, zorder=4)
        else:
            ax.text((xlo + xhi) / 2, 0.5, "no expectation range",
                    ha="center", va="center", fontsize=8.5, color=MPL_GREY,
                    style="italic")
            ax.scatter([actual], [0.5], marker="D", s=95, color=color, alpha=0.5,
                       edgecolor="0.4", linewidths=0.5, zorder=4)
    axes[0].annotate("pred \u00b1 k\u00b7error", xy=(pred + half, 0.72),
                     xytext=(pred + half + 0.06, 0.95), fontsize=8,
                     color=MPL_BLUE,
                     arrowprops=dict(arrowstyle="-", color=MPL_BLUE, lw=0.8))
    axes[-1].set_xlabel("consensus daily rainfall (in)", fontsize=8.5)
    axes[-1].tick_params(labelsize=7.5)
    fig.subplots_adjust(left=0.30, right=0.98, top=0.92, bottom=0.16, hspace=0.25)
    fig.savefig(out_path, transparent=True, pad_inches=0.02)
    plt.close(fig)


# --------------------------------------------------------------------------- #
# python-pptx helpers (shared style with make_qc1_slide_pptx.py)
# --------------------------------------------------------------------------- #
def _style_text(shape, text: str, *, size: float, color: RGBColor,
                bold: bool) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_card(slide, left, top, width, height, text, *, fill, line, ink,
             size=16, bold=True, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.5)
    shp.shadow.inherit = False
    _style_text(shp, text, size=size, color=ink, bold=bold)
    return shp


def add_arrow(slide, left, top, width, height):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = ARROW_FILL
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_caption(slide, left, top, width, text, *, size=12, color=INK, bold=False,
                height=0.34):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    _style_text(box, text, size=size, color=color, bold=bold)
    return box


def add_picture(slide, path, left, top, width, height):
    return slide.shapes.add_picture(
        str(path), Inches(left), Inches(top), Inches(width), Inches(height)
    )


def add_stats_table(slide, left, top, width, height, data, *, header_fill):
    rows, cols = len(data), len(data[0])
    gshape = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = gshape.table
    table.first_row = False
    table.horz_banding = False
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else WHITE
            cell.text = str(data[r][c])
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0]
            run.font.size = Pt(13 if r == 0 else 14)
            run.font.bold = (r == 0 or c == 0)
            run.font.color.rgb = WHITE if r == 0 else INK
    return table


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title(slide, text):
    tb = slide.shapes.add_textbox(Inches(0), Inches(0.12), SLIDE_W, Inches(0.55))
    _style_text(tb, text, size=22, color=INK, bold=True)


# --------------------------------------------------------------------------- #
# Slides
# --------------------------------------------------------------------------- #
def build_slide1(prs, tmp: Path) -> None:
    slide = _blank(prs)
    _title(slide, "QC2 \u00b7 Stage 1 \u2014 regional neighbour statistics")

    m = tmp / "neigh_map.png"
    render_neighbour_map(m)
    add_picture(slide, m, 0.55, 1.15, 5.3, 5.3)
    add_caption(slide, 0.55, 6.55, 5.3,
                "QC1-pass neighbours within 20 km & 50 km, same calendar day",
                size=13, bold=True)

    add_arrow(slide, 6.15, 3.60, 0.85, 0.5)
    add_caption(slide, 6.05, 3.10, 1.05, "median", size=12)

    add_caption(slide, 7.35, 1.55, 5.3, "Statistics for this station-day",
                size=15, bold=True)
    data = [
        ["ring", "n", "median", "MAD"],
        ["20 km", "6", "0.24 in", "0.05"],
        ["50 km", "18", "0.21 in", "0.07"],
    ]
    add_stats_table(slide, 7.35, 2.05, 5.35, 1.65, data, header_fill=BLUE_LINE)

    add_card(slide, 7.35, 4.25, 5.35, 1.55,
             "Robust (median / MAD) so a single bad neighbour cannot distort it.\n"
             "Only QC1-pass neighbours count; the target is excluded.\n"
             "Computed for every located station-day.",
             fill=GREY_FILL, line=GREY_LINE, ink=INK, size=13, bold=False)


def build_slide2(prs, tmp: Path) -> None:
    slide = _blank(prs)
    _title(slide, "QC2 \u00b7 Stage 2 \u2014 expectation models fitted to QC1-pass days")

    t = tmp / "train_table.png"
    render_training_table(t)
    add_picture(slide, t, 0.45, 1.45, 4.4, 4.2)
    add_caption(slide, 0.45, 5.75, 4.4, "Reliable training rows (QC1 pass)",
                size=13, bold=True)

    add_arrow(slide, 5.00, 3.35, 0.6, 0.44)

    tree = tmp / "trees.png"
    render_tree_icon(tree)
    add_picture(slide, tree, 5.85, 1.55, 2.7, 1.15)
    add_card(slide, 5.75, 2.85, 2.85, 0.95,
             "Model 1 (XGBoost)\nstats \u2192 expected rainfall",
             fill=GREY_FILL, line=GREY_LINE, ink=INK, size=12.5)
    add_card(slide, 5.75, 3.95, 2.85, 0.95,
             "Model 2 (XGBoost)\nstats \u2192 expected error",
             fill=BLUE_FILL, line=BLUE_LINE, ink=BLUE_INK, size=12.5)
    add_caption(slide, 5.75, 5.00, 2.85, "learned in log(1+x) space", size=11)

    add_arrow(slide, 8.80, 3.35, 0.6, 0.44)

    c = tmp / "calib.png"
    render_calibration_scatter(c)
    add_picture(slide, c, 9.55, 1.65, 3.4, 3.4)
    add_caption(slide, 9.55, 5.15, 3.4,
                f"k set so {int(COVERAGE_TARGET*100)}% fall in pred \u00b1 k\u00b7error",
                size=12.5, bold=True)


def build_slide3(prs, tmp: Path) -> None:
    slide = _blank(prs)
    _title(slide, "QC2 \u00b7 Stage 2 \u2014 re-judging the days QC1 failed")

    add_card(slide, 0.6, 1.75, 2.3, 1.05, "A day QC1\nfailed",
             fill=RED_FILL, line=RED_LINE, ink=RED_INK, size=14)
    add_arrow(slide, 3.0, 2.05, 0.55, 0.44)
    add_card(slide, 3.65, 1.75, 2.55, 1.05, "Regional stats\n+ Models 1 & 2",
             fill=GREY_FILL, line=GREY_LINE, ink=INK, size=13)
    add_arrow(slide, 6.35, 2.05, 0.55, 0.44)
    add_card(slide, 7.0, 1.6, 3.3, 1.35, "Expectation range\npred \u00b1 k\u00b7error",
             fill=BLUE_FILL, line=BLUE_LINE, ink=BLUE_INK, size=15)

    o = tmp / "outcomes.png"
    render_outcomes(o)
    add_picture(slide, o, 1.05, 3.35, 11.2, 3.15)
    add_caption(slide, 1.05, 6.6, 11.2,
                "inside \u2192 pass (rescued)      outside \u2192 fail (suspect)      "
                "no neighbours \u2192 indeterminate",
                size=13, bold=True)


# --------------------------------------------------------------------------- #
# Deck assembly
# --------------------------------------------------------------------------- #
def build_deck(output_path: Path = OUTPUT_PATH) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    with tempfile.TemporaryDirectory(dir="/var/tmp") as td:
        tmp = Path(td)
        build_slide1(prs, tmp)
        build_slide2(prs, tmp)
        build_slide3(prs, tmp)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

    return output_path


def main() -> None:
    path = build_deck()
    print(f"Wrote QC2 process deck (3 slides) -> {path}")


if __name__ == "__main__":
    main()
